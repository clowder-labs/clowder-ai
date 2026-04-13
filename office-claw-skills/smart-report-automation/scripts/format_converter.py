#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
多格式批量转换器
功能：Word/Excel/PPT ↔ PDF/CSV/图片 批量转换
"""

import argparse
import json
import os
import sys
from pathlib import Path
from typing import List, Tuple, Optional

try:
    from PIL import Image
except ImportError:
    print("错误：缺少 pillow 库，请运行: pip install pillow>=10.0.0")
    sys.exit(1)


class FormatConverter:
    """格式转换器"""
    
    # 支持的转换矩阵
    CONVERSION_MATRIX = {
        "docx": ["pdf", "png", "jpg"],
        "doc": ["pdf", "png", "jpg"],
        "xlsx": ["pdf", "csv"],
        "xls": ["pdf", "csv"],
        "pptx": ["pdf", "png", "jpg"],
        "ppt": ["pdf", "png", "jpg"],
        "pdf": ["png", "jpg"],
    }
    
    def __init__(self, output_dir: str, preserve_format: bool = True):
        """
        初始化转换器
        
        Args:
            output_dir: 输出目录
            preserve_format: 是否保留格式
        """
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.preserve_format = preserve_format
        self.conversion_method = self._detect_conversion_method()
    
    def _detect_conversion_method(self) -> str:
        """
        检测可用的转换方法
        
        Returns:
            转换方法: "libreoffice", "win32com", "python"
        """
        # 优先检测 LibreOffice
        libreoffice_paths = [
            "/usr/bin/libreoffice",
            "/usr/bin/soffice",
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",
            "C:\\Program Files\\LibreOffice\\program\\soffice.exe",
            "C:\\Program Files (x86)\\LibreOffice\\program\\soffice.exe",
        ]
        
        for path in libreoffice_paths:
            if os.path.exists(path):
                return "libreoffice"
        
        # Windows 检测 Office COM
        if sys.platform == "win32":
            try:
                import win32com.client
                return "win32com"
            except ImportError:
                pass
        
        # 使用 Python 库作为备选
        return "python"
    
    def convert_word_to_pdf(self, input_file: Path) -> Path:
        """
        Word 转 PDF
        
        Args:
            input_file: 输入 Word 文件
        
        Returns:
            输出 PDF 文件路径
        """
        output_file = self.output_dir / f"{input_file.stem}.pdf"
        
        if self.conversion_method == "libreoffice":
            import subprocess
            cmd = [
                "libreoffice", "--headless", "--convert-to", "pdf",
                "--outdir", str(self.output_dir), str(input_file)
            ]
            subprocess.run(cmd, check=True, capture_output=True)
        
        elif self.conversion_method == "win32com":
            import win32com.client
            word = win32com.client.Dispatch("Word.Application")
            word.Visible = False
            doc = word.Documents.Open(str(input_file.absolute()))
            doc.SaveAs(str(output_file.absolute()), FileFormat=17)
            doc.Close()
            word.Quit()
        
        else:
            # 使用 python-docx + reportlab 作为备选
            try:
                from docx import Document
                from reportlab.pdfgen import canvas
                from reportlab.lib.pagesizes import letter
                
                doc = Document(str(input_file))
                c = canvas.Canvas(str(output_file), pagesize=letter)
                width, height = letter
                
                y = height - 40
                for para in doc.paragraphs:
                    if y < 40:
                        c.showPage()
                        y = height - 40
                    c.drawString(40, y, para.text[:100])
                    y -= 15
                
                c.save()
            except ImportError:
                raise ImportError("缺少必要的库: python-docx, reportlab")
        
        return output_file
    
    def convert_excel_to_pdf(self, input_file: Path) -> Path:
        """
        Excel 转 PDF
        
        Args:
            input_file: 输入 Excel 文件
        
        Returns:
            输出 PDF 文件路径
        """
        output_file = self.output_dir / f"{input_file.stem}.pdf"
        
        if self.conversion_method == "libreoffice":
            import subprocess
            cmd = [
                "libreoffice", "--headless", "--convert-to", "pdf",
                "--outdir", str(self.output_dir), str(input_file)
            ]
            subprocess.run(cmd, check=True, capture_output=True)
        
        elif self.conversion_method == "win32com":
            import win32com.client
            excel = win32com.client.Dispatch("Excel.Application")
            excel.Visible = False
            wb = excel.Workbooks.Open(str(input_file.absolute()))
            wb.SaveAs(str(output_file.absolute()), FileFormat=57)
            wb.Close()
            excel.Quit()
        
        else:
            # 使用 openpyxl 作为备选（简化版本）
            try:
                import openpyxl
                from reportlab.pdfgen import canvas
                from reportlab.lib.pagesizes import letter
                
                wb = openpyxl.load_workbook(str(input_file))
                ws = wb.active
                
                c = canvas.Canvas(str(output_file), pagesize=letter)
                width, height = letter
                
                y = height - 40
                for row in ws.iter_rows(max_row=50):
                    row_text = " | ".join([str(cell.value or "") for cell in row])
                    if y < 40:
                        c.showPage()
                        y = height - 40
                    c.drawString(40, y, row_text[:120])
                    y -= 15
                
                c.save()
            except ImportError:
                raise ImportError("缺少必要的库: openpyxl, reportlab")
        
        return output_file
    
    def convert_excel_to_csv(self, input_file: Path) -> Path:
        """
        Excel 转 CSV
        
        Args:
            input_file: 输入 Excel 文件
        
        Returns:
            输出 CSV 文件路径
        """
        import pandas as pd
        
        output_file = self.output_dir / f"{input_file.stem}.csv"
        df = pd.read_excel(input_file, sheet_name=0)
        df.to_csv(output_file, index=False, encoding="utf-8-sig")
        
        return output_file
    
    def convert_ppt_to_pdf(self, input_file: Path) -> Path:
        """
        PPT 转 PDF
        
        Args:
            input_file: 输入 PPT 文件
        
        Returns:
            输出 PDF 文件路径
        """
        output_file = self.output_dir / f"{input_file.stem}.pdf"
        
        if self.conversion_method == "libreoffice":
            import subprocess
            cmd = [
                "libreoffice", "--headless", "--convert-to", "pdf",
                "--outdir", str(self.output_dir), str(input_file)
            ]
            subprocess.run(cmd, check=True, capture_output=True)
        
        elif self.conversion_method == "win32com":
            import win32com.client
            ppt = win32com.client.Dispatch("PowerPoint.Application")
            presentation = ppt.Presentations.Open(str(input_file.absolute()))
            presentation.SaveAs(str(output_file.absolute()), FileFormat=32)
            presentation.Close()
            ppt.Quit()
        
        else:
            # 使用 python-pptx 作为备选（简化版本）
            try:
                from pptx import Presentation
                from reportlab.pdfgen import canvas
                from reportlab.lib.pagesizes import letter
                
                prs = Presentation(str(input_file))
                c = canvas.Canvas(str(output_file), pagesize=letter)
                width, height = letter
                
                for slide in prs.slides:
                    y = height - 40
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            c.drawString(40, y, shape.text[:100])
                            y -= 15
                    c.showPage()
                
                c.save()
            except ImportError:
                raise ImportError("缺少必要的库: python-pptx, reportlab")
        
        return output_file
    
    def convert_ppt_to_images(self, input_file: Path, image_format: str = "png") -> List[Path]:
        """
        PPT 转图片（每页一张）
        
        Args:
            input_file: 输入 PPT 文件
            image_format: 图片格式（png/jpg）
        
        Returns:
            输出图片文件路径列表
        """
        output_files = []
        
        if self.conversion_method == "libreoffice":
            # 先转 PDF
            pdf_file = self.convert_ppt_to_pdf(input_file)
            
            # PDF 转图片
            try:
                from pdf2image import convert_from_path
                images = convert_from_path(str(pdf_file))
                
                for i, image in enumerate(images, start=1):
                    output_file = self.output_dir / f"{input_file.stem}_page{i}.{image_format}"
                    image.save(output_file, image_format.upper())
                    output_files.append(output_file)
            except ImportError:
                # 使用 PIL 的备选方案
                print("警告：缺少 pdf2image 库，无法转换 PPT 为图片")
        
        elif self.conversion_method == "win32com":
            import win32com.client
            ppt = win32com.client.Dispatch("PowerPoint.Application")
            presentation = ppt.Presentations.Open(str(input_file.absolute()))
            
            # 导出为图片
            export_path = str(self.output_dir / input_file.stem)
            presentation.Export(export_path, image_format)
            presentation.Close()
            ppt.Quit()
            
            # 收集生成的图片
            for file in self.output_dir.glob(f"{input_file.stem}/*.{image_format}"):
                output_files.append(file)
        
        return output_files
    
    def convert_file(
        self,
        input_file: str,
        output_format: str
    ) -> Tuple[bool, str, Optional[str]]:
        """
        转换单个文件
        
        Args:
            input_file: 输入文件路径
            output_format: 目标格式
        
        Returns:
            (是否成功, 输出路径, 错误信息)
        """
        input_path = Path(input_file)
        
        if not input_path.exists():
            return False, "", f"文件不存在: {input_file}"
        
        # 获取文件扩展名
        input_ext = input_path.suffix.lower().lstrip(".")
        
        # 检查是否支持转换
        if input_ext not in self.CONVERSION_MATRIX:
            return False, "", f"不支持的输入格式: {input_ext}"
        
        if output_format.lower() not in self.CONVERSION_MATRIX[input_ext]:
            return False, "", f"不支持从 {input_ext} 转换到 {output_format}"
        
        try:
            output_format = output_format.lower()
            
            # Word 文件转换
            if input_ext in ["docx", "doc"]:
                if output_format == "pdf":
                    output = self.convert_word_to_pdf(input_path)
                else:
                    return False, "", f"暂不支持 Word 转 {output_format}"
            
            # Excel 文件转换
            elif input_ext in ["xlsx", "xls"]:
                if output_format == "pdf":
                    output = self.convert_excel_to_pdf(input_path)
                elif output_format == "csv":
                    output = self.convert_excel_to_csv(input_path)
                else:
                    return False, "", f"暂不支持 Excel 转 {output_format}"
            
            # PPT 文件转换
            elif input_ext in ["pptx", "ppt"]:
                if output_format == "pdf":
                    output = self.convert_ppt_to_pdf(input_path)
                elif output_format in ["png", "jpg"]:
                    outputs = self.convert_ppt_to_images(input_path, output_format)
                    output = str(outputs)
                else:
                    return False, "", f"暂不支持 PPT 转 {output_format}"
            
            else:
                return False, "", f"未实现的转换: {input_ext} -> {output_format}"
            
            return True, str(output), None
        
        except Exception as e:
            return False, "", f"转换失败: {str(e)}"
    
    def batch_convert(
        self,
        input_files: List[str],
        output_format: str
    ) -> List[Tuple[str, bool, str, Optional[str]]]:
        """
        批量转换文件
        
        Args:
            input_files: 输入文件列表
            output_format: 目标格式
        
        Returns:
            转换结果列表: [(输入文件, 是否成功, 输出路径, 错误信息)]
        """
        results = []
        
        for input_file in input_files:
            success, output, error = self.convert_file(input_file, output_format)
            results.append((input_file, success, output, error))
        
        return results


def batch_convert_files(
    input_files: List[str],
    output_format: str,
    output_dir: str = "./output",
    preserve_format: bool = True
) -> List[str]:
    """
    批量转换文件主函数
    
    Args:
        input_files: 输入文件列表
        output_format: 目标格式（pdf/png/jpg/csv）
        output_dir: 输出目录
        preserve_format: 是否保留格式
    
    Returns:
        成功转换的输出文件列表
    """
    converter = FormatConverter(output_dir, preserve_format)
    results = converter.batch_convert(input_files, output_format)
    
    output_paths = []
    for input_file, success, output, error in results:
        if success:
            print(f"✓ {input_file} -> {output}")
            output_paths.append(output)
        else:
            print(f"✗ {input_file}: {error}")
    
    return output_paths


def main():
    """命令行入口"""
    parser = argparse.ArgumentParser(description="多格式批量转换器")
    parser.add_argument("--input", nargs="+", required=True, help="输入文件路径（支持多个）")
    parser.add_argument("--format", required=True, choices=["pdf", "png", "jpg", "csv"], help="目标格式")
    parser.add_argument("--output-dir", default="./output", help="输出目录")
    parser.add_argument("--preserve-format", action="store_true", help="保留格式")
    
    args = parser.parse_args()
    
    results = batch_convert_files(
        input_files=args.input,
        output_format=args.format,
        output_dir=args.output_dir,
        preserve_format=args.preserve_format
    )
    
    print(f"\n转换完成: {len(results)} 个文件")
    return results


if __name__ == "__main__":
    main()
